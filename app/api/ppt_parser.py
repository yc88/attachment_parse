import os
from io import BytesIO
import requests

from app.schemas.const import MyErrorCode, HttpStatusCode
from app.schemas.parse import FileParseResult
from app.utils.app_exceptions import UnicornException
from app.api.content_regex import RegexRules, get_regex_val
from app.schemas.parse import RegexField, RegexFieldStatistical, ExtraField
from datetime import datetime, timezone, timedelta
from pptx import Presentation
from PIL import Image
from app.api.img_parser import get_return_result as get_return_img_result


def parse_local_ppt_file(file_path, tesseract_path: str, lang: str = None):
    """
    解析本地的ppt文件内容
    """
    try:
        encoding = 'utf-8'
        # 获取文件名称
        file_name = os.path.basename(file_path)
        # 获取文件大小
        file_size = os.path.getsize(file_path)
        # 获取最后修改时间
        last_modified = datetime.fromtimestamp(os.path.getmtime(file_path)).strftime('%Y-%m-%d %H:%M:%S')
        py_bytes = BytesIO(open(file_path, 'rb').read())
        return get_return_result(py_bytes, encoding, file_name, file_size, last_modified, tesseract_path, lang)
    except Exception as e:
        raise UnicornException(HttpStatusCode.HTTP_500_INTERNAL_SERVER_ERROR,
                               MyErrorCode.InvalidArgument.value,
                               e.__str__())


def parse_remote_ppt_file(file_url, tesseract_path: str, lang: str = None):
    try:
        encoding = 'utf-8'
        # 下载远程文件内容
        response = requests.get(file_url)
        if response.status_code == 200:
            # 获取文件名称
            file_name = os.path.basename(file_url)
            # 文件内容
            raw_content = response.content
            # 获取文件大小
            file_size = len(raw_content)

            # 获取最后修改时间
            last_modified = response.headers.get('Last-Modified')
            if last_modified:
                last_modified_datetime = datetime.strptime(last_modified,
                                                           '%a, %d %b %Y %H:%M:%S %Z').replace(tzinfo=timezone.utc)
                cst_time = last_modified_datetime.astimezone(timezone(timedelta(hours=8)))
                last_modified = cst_time.strftime('%Y-%m-%d %H:%M:%S')
            py_bytes = BytesIO(raw_content)
            return get_return_result(py_bytes, encoding, file_name, file_size, last_modified, tesseract_path, lang)
        else:
            raise UnicornException(HttpStatusCode.HTTP_500_INTERNAL_SERVER_ERROR,
                                   MyErrorCode.InvalidArgument.value,
                                   " 请求错误 ")
    except Exception as e:
        raise UnicornException(HttpStatusCode.HTTP_500_INTERNAL_SERVER_ERROR,
                               MyErrorCode.InvalidArgument.value,
                               e.__str__())


def get_return_result(py_bytes, encoding: str, file_name: str, file_size: int,
                      last_modified: str, tesseract_path: str, lang: str = None) -> FileParseResult:
    """ return the same result """
    prs = Presentation(py_bytes)
    content = []
    title_core = prs.core_properties.title
    title = []
    images_txt = []
    for slide in prs.slides:
        shapes_title = slide.shapes.title
        if shapes_title and shapes_title.text:
            title.append(shapes_title.text.strip())
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                content.append(shape.text)
            if shape.shape_type == 19:  # Table shape
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            content.append(cell.text.strip())
            elif shape.shape_type == 13:  # Picture shape
                image_bytes = shape.image.blob
                image = Image.open(BytesIO(image_bytes))
                images_txt.append(get_return_img_result(image, encoding, '', 0, '', tesseract_path, lang).file_content)

    file_content = '\n'.join(content) + '\n'.join(images_txt)
    regex_all, regex_all_statistical = get_regex_val(file_content)
    if not title:
        title.append(title_core)
    extra_field = ExtraField(
        date='',
        language='',
        charset='',
        subject='\n'.join(title),
        from_email='',
        to_email='',
        images=[],
        link=[],
    )
    return FileParseResult(file_size=file_size, file_name=file_name, last_modified=last_modified,
                           file_content=file_content, file_encoding=encoding,
                           regex_all=regex_all,
                           regex_all_statistical=regex_all_statistical, extra=extra_field)
